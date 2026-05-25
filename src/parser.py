# -*- coding: utf-8 -*-
"""Multi-format document parser (PDF/PPTX/DOCX/TXT/SRT)."""

from __future__ import annotations

from pathlib import Path


def parse_file(file_path: Path) -> str:
    suffix = file_path.suffix.lower()
    try:
        if suffix in (".txt", ".srt", ".vtt"):
            return _txt(file_path)
        if suffix == ".pdf":
            return _pdf(file_path)
        if suffix == ".pptx":
            return _pptx(file_path)
        if suffix == ".docx":
            return _docx(file_path)
    except Exception:
        return ""
    return ""


def _txt(path: Path) -> str:
    for encoding in ("utf-8", "utf-8-sig", "cp950", "big5", "latin1"):
        try:
            return path.read_text(encoding=encoding).strip()
        except (UnicodeDecodeError, ValueError):
            continue
    return ""


def _pdf(path: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    pages = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            pages.append(text.strip())
    return "\n".join(pages)


def _pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = paragraph.text.strip()
                    if line:
                        texts.append(line)
    return "\n".join(texts)


def _docx(path: Path) -> str:
    from docx import Document

    doc = Document(str(path))
    lines = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    return "\n".join(lines)
